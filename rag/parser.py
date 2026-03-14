"""
文档解析模块 —— 离线预处理阶段使用

将各种格式的本地文档转换为纯文本字符串，供 ingest.py 后续分块和向量化使用。
支持格式：.md / .txt / .html / .pdf / .docx / .pptx / .xlsx

注意：网页 URL 的实时抓取由 agent/tools.py 中的 fetch_url 工具负责，
      本模块只处理本地文件。
"""

from pathlib import Path


# 支持的文件扩展名集合
SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(
    {".md", ".txt", ".html", ".htm", ".pdf", ".docx", ".pptx", ".xlsx"}
)


def parse_file(file_path: str | Path) -> str:
    """
    解析单个本地文件，返回纯文本字符串。

    Args:
        file_path: 文件路径（支持字符串或 Path 对象）。

    Returns:
        文件的纯文本内容，末尾不带多余空白。

    Raises:
        ValueError: 文件格式不受支持。
        FileNotFoundError: 文件不存在。
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {path}")

    suffix = path.suffix.lower()

    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"不支持的文件格式: '{suffix}'，"
            f"支持的格式: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    match suffix:
        case ".md" | ".txt":
            return _parse_text(path)
        case ".html" | ".htm":
            return _parse_html(path)
        case ".pdf":
            return _parse_pdf(path)
        case ".docx":
            return _parse_docx(path)
        case ".pptx":
            return _parse_pptx(path)
        case ".xlsx":
            return _parse_xlsx(path)
        case _:
            # 理论上不会到这里，frozenset 已经过滤
            raise ValueError(f"未处理的格式: '{suffix}'")


# ── 各格式解析实现 ────────────────────────────────────────────────────────────


def _parse_text(path: Path) -> str:
    """解析 .md / .txt：直接读取文本，自动检测编码。"""
    # 优先 utf-8，失败则回退 gbk（兼容 Windows 中文环境）
    for encoding in ("utf-8", "gbk", "latin-1"):
        try:
            return path.read_text(encoding=encoding).strip()
        except UnicodeDecodeError:
            continue
    # latin-1 是超集，理论上不会失败，此处保险起见
    return path.read_text(errors="replace").strip()


def _parse_html(path: Path) -> str:
    """解析 .html / .htm：用 BeautifulSoup 提取正文，去除脚本和样式标签。"""
    from bs4 import BeautifulSoup

    html = _parse_text(path)
    soup = BeautifulSoup(html, "lxml")

    # 移除不需要的标签
    for tag in soup(["script", "style", "head", "nav", "footer", "aside"]):
        tag.decompose()

    # 提取纯文本，保留段落间空行
    lines = (line.strip() for line in soup.get_text(separator="\n").splitlines())
    # 过滤空行连续出现（最多保留一个空行）
    result: list[str] = []
    prev_blank = False
    for line in lines:
        if line:
            result.append(line)
            prev_blank = False
        elif not prev_blank:
            result.append("")
            prev_blank = True

    return "\n".join(result).strip()


def _parse_pdf(path: Path) -> str:
    """解析 .pdf：用 pypdf 逐页提取文本，页间用换行分隔。"""
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        text = text.strip()
        if text:
            pages.append(text)

    return "\n\n".join(pages).strip()


def _parse_docx(path: Path) -> str:
    """解析 .docx：提取所有段落文本，保留段落换行。"""
    from docx import Document

    doc = Document(str(path))
    paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
    return "\n\n".join(paragraphs).strip()


def _parse_pptx(path: Path) -> str:
    """解析 .pptx：遍历所有 slide 的所有 shape，提取文本框内容。"""
    from pptx import Presentation

    prs = Presentation(str(path))
    slides_text: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:  # type: ignore[union-attr]
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            slides_text.append(f"[Slide {slide_idx}]\n" + "\n".join(texts))

    return "\n\n".join(slides_text).strip()


def _parse_xlsx(path: Path) -> str:
    """
    解析 .xlsx：遍历所有 sheet 的所有行，
    每行转为 '列1 | 列2 | 列3 ...' 格式，便于语义检索。
    """
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sheets_text: list[str] = []

    for sheet in wb.worksheets:
        rows_text: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() if cell is not None else "" for cell in row]
            # 跳过全空行
            if any(cells):
                rows_text.append(" | ".join(cells))
        if rows_text:
            sheets_text.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows_text))

    wb.close()
    return "\n\n".join(sheets_text).strip()
