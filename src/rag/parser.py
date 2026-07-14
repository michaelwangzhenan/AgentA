"""
文档解析模块 —— 离线预处理阶段使用

将各种格式的本地文档转换为纯文本字符串，供 ingest.py 后续分块和向量化使用。
支持格式：.md / .txt / .html / .pdf / .docx / .pptx / .xlsx

注意：网页 URL 的实时抓取由 agent/tools.py 中的 fetch_url 工具负责，
      本模块只处理本地文件。
"""

import logging
import re
import subprocess
import sys
import tempfile
import zipfile
from collections import Counter
from collections.abc import Iterator
from contextlib import contextmanager
from pathlib import Path

logger = logging.getLogger(__name__)


# 支持的文件扩展名集合
SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(
    {".md", ".txt", ".html", ".htm", ".pdf", ".docx", ".pptx", ".xlsx"}
)


class DocxParseError(ValueError):
    """DOCX 安全检查或隔离解析失败。"""


def is_office_temp_file(path_or_name: str | Path) -> bool:
    """判断是否为 Office 自动生成的 ~$ 临时文件。"""
    return Path(path_or_name).name.startswith("~$")


def inspect_docx_uncompressed_size(path: Path, max_bytes: int) -> int:
    """返回 DOCX 解压总量；超过上限时在加载 XML 前拒绝。"""
    total = 0
    try:
        with zipfile.ZipFile(path) as archive:
            for info in archive.infolist():
                total += info.file_size
                if total > max_bytes:
                    raise DocxParseError(
                        "DOCX 解压后过大"
                        f"（{total / 1024 / 1024:.1f} MiB > {max_bytes / 1024 / 1024:.0f} MiB）"
                    )
    except zipfile.BadZipFile as exc:
        raise DocxParseError(f"DOCX 文件损坏或格式无效: {path.name}") from exc
    return total


@contextmanager
def parsed_docx_temp(path: Path) -> Iterator[Path]:
    """在受限子进程中解析 DOCX，返回只含清洗文本的临时文件。"""
    cfg = _cfg()
    max_bytes = cfg.DOCX_MAX_UNZIP_MB * 1024 * 1024
    uncompressed = inspect_docx_uncompressed_size(path, max_bytes)
    logger.info(
        "[parser] DOCX 预检通过: %s compressed=%d uncompressed=%d",
        path.name,
        path.stat().st_size,
        uncompressed,
    )

    temp = tempfile.NamedTemporaryFile(prefix="agenta-docx-", suffix=".txt", delete=False)
    temp_path = Path(temp.name)
    temp.close()
    command = [
        sys.executable,
        "-m",
        "src.rag.docx_worker",
        str(path),
        "--memory-mb",
        str(cfg.DOCX_PARSE_MEMORY_MB),
    ]
    project_root = Path(__file__).resolve().parents[2]
    try:
        with temp_path.open("wb") as output:
            process = subprocess.Popen(
                command,
                cwd=project_root,
                stdout=output,
                stderr=subprocess.PIPE,
            )
            try:
                _, stderr = process.communicate(timeout=cfg.DOCX_PARSE_TIMEOUT_SEC)
            except subprocess.TimeoutExpired as exc:
                process.kill()
                process.communicate()
                raise DocxParseError(
                    f"DOCX 解析超时（超过 {cfg.DOCX_PARSE_TIMEOUT_SEC} 秒）: {path.name}"
                ) from exc

        if process.returncode != 0:
            detail = (stderr or b"").decode("utf-8", errors="replace").strip()
            if process.returncode in (-9, 137):
                detail = detail or f"超过 {cfg.DOCX_PARSE_MEMORY_MB} MiB 内存限制"
            raise DocxParseError(
                f"DOCX 隔离解析失败: {path.name}"
                + (f" — {detail[-500:]}" if detail else "")
            )
        logger.info("[parser] DOCX 隔离解析完成: %s text_bytes=%d", path.name, temp_path.stat().st_size)
        yield temp_path
    finally:
        temp_path.unlink(missing_ok=True)


# ── 解析后内容清洗 ────────────────────────────────────────────────────────────
# 目的：去除 PDF/Word/PPT 等长文档常见的"模板噪声"——页眉页脚、版权声明、纯页码、
# PPT 母版水印等。这些片段会大量重复进入向量库，污染相似度计算并稀释命中率。

# 满足以下任一条件即被视为"模板/页脚噪声"：
#   1. 整行只有数字、日期、版权字符（©/® 等）
#   2. 整行匹配典型页码模式（"1 / 32" 或 "Page 5 of 10"）
#   3. 整行匹配典型版权/保留声明（不限语种）
_NOISE_LINE_PATTERNS: tuple[re.Pattern[str], ...] = (
    re.compile(r"^\s*\d+\s*$"),                                # 纯页码
    re.compile(r"^\s*[-–—•·]+\s*\d+\s*[-–—•·]+\s*$"),          # "- 12 -" / "— 12 —"
    re.compile(r"^\s*page\s+\d+\s*(of\s+\d+)?\s*$", re.I),      # Page 5 / Page 5 of 10
    re.compile(r"^\s*\d+\s*[/／]\s*\d+\s*$"),                    # 1/32 / 1／32
    re.compile(r"©|®|™"),                                       # 含版权符号
    re.compile(r"copyright|all\s+rights\s+reserved", re.I),     # 英文版权
    re.compile(r"版权所有|保留所有权利|未经.*授权.*禁止"),      # 中文版权
)

# 短行重复阈值：一份文档中重复出现 ≥ N 次的短行（≤ MAX_LEN 字符）视为页眉页脚/水印
_NOISE_REPEAT_THRESHOLD: int = 5
_NOISE_MAX_LEN: int = 80


def _is_noise_line(line: str) -> bool:
    """单行级别的噪声检测（不需要全文上下文，纯规则）。"""
    s = line.strip()
    if not s:
        return False
    return any(p.search(s) for p in _NOISE_LINE_PATTERNS)


def _clean_extracted_text(text: str) -> str:
    """
    解析层最后一步：去除模板/页脚类噪声，避免污染向量库。

    清洗策略：
      1. 整行规则匹配：纯页码、页码模式、版权声明 → 直接丢弃。
      2. 跨页重复：在长文档中重复出现 ≥ _NOISE_REPEAT_THRESHOLD 次的"短行"
         （length ≤ _NOISE_MAX_LEN）视为页眉页脚/水印 → 丢弃。
      3. 多余空行折叠为最多一个空行（避免 chunk 中出现大段空白）。
    """
    if not text or not text.strip():
        return ""

    raw_lines = text.splitlines()

    # Pass 1: 统计短行重复次数
    short_line_counts: Counter[str] = Counter()
    for line in raw_lines:
        s = line.strip()
        if 0 < len(s) <= _NOISE_MAX_LEN:
            short_line_counts[s] += 1

    repeated_noise: set[str] = {
        s for s, n in short_line_counts.items() if n >= _NOISE_REPEAT_THRESHOLD
    }

    # Pass 2: 应用规则 + 重复短行 + 空行折叠
    cleaned: list[str] = []
    prev_blank = False
    dropped = 0
    for line in raw_lines:
        stripped = line.strip()
        if not stripped:
            if not prev_blank:
                cleaned.append("")
                prev_blank = True
            continue
        if _is_noise_line(stripped) or stripped in repeated_noise:
            dropped += 1
            continue
        cleaned.append(stripped)
        prev_blank = False

    if dropped > 0:
        logger.debug("[parser] 清洗丢弃噪声行 %d 条", dropped)

    return "\n".join(cleaned).strip()


def parse_file(file_path: str | Path) -> str:
    """
    解析单个本地文件，返回纯文本字符串。

    Args:
        file_path: 文件路径（支持字符串或 Path 对象）。

    Returns:
        文件的纯文本内容，末尾不带多余空白。

    Raises:
        ValueError: 文件格式不受支持。
        FileNotFoundError: 文件不存在。
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {path}")

    if is_office_temp_file(path):
        raise ValueError(f"跳过 Office 临时文件: {path.name}")

    suffix = path.suffix.lower()

    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"不支持的文件格式: '{suffix}'，"
            f"支持的格式: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    match suffix:
        case ".md" | ".txt":
            raw = _parse_text(path)
        case ".html" | ".htm":
            raw = _parse_html(path)
        case ".pdf":
            raw = _parse_pdf(path)
        case ".docx":
            with parsed_docx_temp(path) as parsed_path:
                return parsed_path.read_text(encoding="utf-8")
        case ".pptx":
            raw = _parse_pptx(path)
        case ".xlsx":
            raw = _parse_xlsx(path)
        case _:
            # 理论上不会到这里，frozenset 已经过滤
            raise ValueError(f"未处理的格式: '{suffix}'")

    # 统一过一遍模板/页脚噪声清洗，避免污染下游 chunk → embedding
    return _clean_extracted_text(raw)


# ── 各格式解析实现 ────────────────────────────────────────────────────────────


def _parse_text(path: Path) -> str:
    """解析 .md / .txt：直接读取文本，自动检测编码。"""
    # 优先 utf-8，失败则回退 gbk（兼容 Windows 中文环境）
    for encoding in ("utf-8", "gbk", "latin-1"):
        try:
            return path.read_text(encoding=encoding).strip()
        except UnicodeDecodeError:
            continue
    # latin-1 是超集，理论上不会失败，此处保险起见
    return path.read_text(errors="replace").strip()


def _parse_html(path: Path) -> str:
    """
    解析 .html / .htm：用 BeautifulSoup 提取正文，去除脚本/样式/导航/页脚标签。

    结构保留：将 h1~h6 标签替换为 Markdown 风格 "# 标题" 行，下游 splitter 据此切段
    并注入 heading_path metadata。
    """
    from bs4 import BeautifulSoup, NavigableString

    html = _parse_text(path)
    soup = BeautifulSoup(html, "lxml")

    # 移除不需要的标签
    for tag in soup(["script", "style", "head", "nav", "footer", "aside"]):
        tag.decompose()

    # h1~h6 → Markdown 标题；用 NavigableString 替换确保 get_text 正常提取
    for h_tag in soup.find_all(re.compile(r"^h[1-6]$")):
        try:
            level = int(h_tag.name[1])
        except (TypeError, ValueError):
            continue
        title = h_tag.get_text(" ", strip=True)
        if not title:
            h_tag.decompose()
            continue
        h_tag.replace_with(NavigableString(f"\n{'#' * level} {title}\n"))

    # 提取纯文本，保留段落间空行
    lines = (line.strip() for line in soup.get_text(separator="\n").splitlines())
    # 过滤空行连续出现（最多保留一个空行）
    result: list[str] = []
    prev_blank = False
    for line in lines:
        if line:
            result.append(line)
            prev_blank = False
        elif not prev_blank:
            result.append("")
            prev_blank = True

    return "\n".join(result).strip()


def _parse_pdf(path: Path) -> str:
    """
    解析 .pdf：多 backend 递降 fallback + 扫描版自动 OCR（Iter-4）。

    流程：
      1. 文本层提取按优先级 pymupdf(fitz) → pdfplumber → pypdf 顺序尝试，
         任一可用且产出非空即返回。pymupdf 文本质量与速度显著优于 pypdf，
         pdfplumber 对带表格/分栏的复杂排版处理更精细。
      2. 提取完成后检测"平均每页字符数"，低于 config.RAG_OCR_TRIGGER_CHARS_PER_PAGE
         时认为是扫描版/图片版 PDF，调用 rapidocr-onnxruntime 对全文 OCR。
      3. 三个 PDF 库与 rapidocr-onnxruntime 全部按软依赖处理；缺哪个都不报错，
         仅在日志里提示用户安装。任何一层成功即采用其结果。

    所有页前都会插入 "[[PAGE:N]]" 锚点供 splitter 识别。
    """
    pages = _pdf_extract_text_pages(path)
    if pages:
        total_chars = sum(len(t) for _, t in pages)
        avg = total_chars / max(len(pages), 1)
        # 文本密度过低 → 大概率扫描版；尝试 OCR 兜底
        if (
            _cfg().RAG_OCR_FALLBACK_ENABLED
            and avg < _cfg().RAG_OCR_TRIGGER_CHARS_PER_PAGE
        ):
            logger.info(
                "[parser] PDF 文本密度极低（%.1f 字符/页 < 阈值 %d），尝试 OCR fallback",
                avg, _cfg().RAG_OCR_TRIGGER_CHARS_PER_PAGE,
            )
            ocr_pages = _pdf_ocr_pages(path)
            if ocr_pages:
                pages = ocr_pages
    elif _cfg().RAG_OCR_FALLBACK_ENABLED:
        # 文本层完全提不出来 → 整本扫描；唯一可行的就是 OCR
        logger.info("[parser] 所有 PDF 后端均无文本，尝试 OCR fallback")
        pages = _pdf_ocr_pages(path)

    if not pages:
        return ""

    return "\n\n".join(f"[[PAGE:{n}]]\n{t.strip()}" for n, t in pages if t.strip()).strip()


def _cfg():
    """延迟拿 src.config，避免循环引用与导入顺序问题。"""
    import src.config as _c
    return _c


def _pdf_extract_text_pages(path: Path) -> list[tuple[int, str]]:
    """按优先级尝试 pymupdf → pdfplumber → pypdf 三个文本提取后端。"""
    for backend in (_pdf_pymupdf, _pdf_pdfplumber, _pdf_pypdf):
        try:
            pages = backend(path)
        except ImportError:
            continue   # 该后端未安装，下一个
        except Exception as e:
            logger.warning("[parser] PDF 后端 %s 解析失败: %s", backend.__name__, e)
            continue
        if pages:
            logger.debug("[parser] PDF 解析使用后端: %s（%d 页）",
                         backend.__name__, len(pages))
            return pages
    return []


def _pdf_pymupdf(path: Path) -> list[tuple[int, str]]:
    """pymupdf (fitz)：最快、文本质量最佳。"""
    import fitz  # PyMuPDF, 软依赖
    pages: list[tuple[int, str]] = []
    doc = fitz.open(str(path))
    try:
        for idx, page in enumerate(doc, start=1):
            text = (page.get_text() or "").strip()
            if text:
                pages.append((idx, text))
    finally:
        doc.close()
    return pages


def _pdf_pdfplumber(path: Path) -> list[tuple[int, str]]:
    """pdfplumber：擅长表格/分栏。"""
    import pdfplumber  # 软依赖
    pages: list[tuple[int, str]] = []
    with pdfplumber.open(str(path)) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            try:
                text = (page.extract_text() or "").strip()
            except Exception:
                text = ""
            if text:
                pages.append((idx, text))
    return pages


def _pdf_pypdf(path: Path) -> list[tuple[int, str]]:
    """pypdf：兜底（项目原依赖，必定可用）。"""
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages: list[tuple[int, str]] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append((idx, text))
    return pages


def _pdf_ocr_pages(path: Path) -> list[tuple[int, str]]:
    """
    rapidocr-onnxruntime + pymupdf 渲染图片做 OCR。

    依赖均为软依赖，缺任一直接返回空列表；外层会保留文本层结果（如有）。
    """
    try:
        from rapidocr_onnxruntime import RapidOCR
    except ImportError:
        logger.info(
            "[parser] rapidocr-onnxruntime 未安装，无法 OCR 扫描版 PDF。"
            "如需启用：pip install rapidocr-onnxruntime"
        )
        return []
    try:
        import fitz  # PyMuPDF 用于把 PDF 渲染成图片
    except ImportError:
        logger.info(
            "[parser] pymupdf 未安装，无法将 PDF 渲染为图片喂给 OCR。"
            "如需启用：pip install pymupdf"
        )
        return []

    cfg = _cfg()
    try:
        ocr = RapidOCR()
    except Exception as e:
        logger.warning("[parser] RapidOCR 初始化失败：%s", e)
        return []

    pages: list[tuple[int, str]] = []
    doc = fitz.open(str(path))
    page_count = 0
    try:
        page_count = len(doc) if hasattr(doc, "__len__") else 0
        for idx, page in enumerate(doc, start=1):
            try:
                pix = page.get_pixmap(dpi=cfg.RAG_OCR_DPI)
                img_bytes = pix.tobytes("png")
                result, _ = ocr(img_bytes)
            except Exception as e:
                logger.warning("[parser] OCR 第 %d 页失败：%s", idx, e)
                continue
            if not result:
                continue
            # rapidocr 返回 [(box, text, score), ...]
            page_text = "\n".join(
                item[1] for item in result if item and len(item) >= 2 and item[1]
            )
            page_text = page_text.strip()
            if page_text:
                pages.append((idx, page_text))
    finally:
        doc.close()

    if pages:
        total = sum(len(t) for _, t in pages)
        logger.info(
            "[parser] OCR 完成：%d 页有效 / 共 %d 页，总字符数 %d",
            len(pages), page_count or len(pages), total,
        )
    return pages


def _parse_docx(path: Path) -> str:
    """
    解析 .docx：按段落顺序提取文本；识别 "Heading 1"~"Heading 9" 样式的段落为
    Markdown 标题（前缀 #~#########），下游 splitter 据此切段并注入 heading_path。
    """
    from docx import Document

    doc = Document(str(path))
    paragraphs: list[str] = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style_name = ""
        try:
            style_name = (para.style.name or "") if para.style else ""
        except Exception:
            style_name = ""
        m = re.match(r"^Heading\s+([1-9])$", style_name)
        if m:
            level = int(m.group(1))
            paragraphs.append(f"{'#' * level} {text}")
        else:
            paragraphs.append(text)
    return "\n\n".join(paragraphs).strip()


def _parse_pptx(path: Path) -> str:
    """
    解析 .pptx：每张 slide 前插入 [[PAGE:N]] 锚点 + "[Slide N]" 字面（向后兼容）；
    若 slide 含 title placeholder，再以 "## <title>" 形式注入 Markdown 标题。
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    slides_text: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # 优先识别 title placeholder
        slide_title = ""
        try:
            title_shape = slide.shapes.title  # type: ignore[union-attr]
            if title_shape is not None and title_shape.has_text_frame:
                slide_title = (title_shape.text_frame.text or "").strip()
        except Exception:
            slide_title = ""

        body_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:  # type: ignore[union-attr]
                line = "".join(run.text for run in para.runs).strip()
                if line and line != slide_title:
                    body_lines.append(line)

        if not (slide_title or body_lines):
            continue

        block_lines: list[str] = [f"[[PAGE:{slide_idx}]]", f"[Slide {slide_idx}]"]
        if slide_title:
            block_lines.append(f"## {slide_title}")
        block_lines.extend(body_lines)
        slides_text.append("\n".join(block_lines))

    return "\n\n".join(slides_text).strip()


def _parse_xlsx(path: Path) -> str:
    """
    解析 .xlsx：每个 sheet 用 "# Sheet: <name>" 作 Markdown 标题（供 splitter 切段），
    同时保留 "[Sheet: <name>]" 字面（向后兼容）；行内单元格用 " | " 分隔。
    """
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sheets_text: list[str] = []

    for sheet in wb.worksheets:
        rows_text: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() if cell is not None else "" for cell in row]
            if any(cells):
                rows_text.append(" | ".join(cells))
        if rows_text:
            block = (
                f"# Sheet: {sheet.title}\n"
                f"[Sheet: {sheet.title}]\n"
                + "\n".join(rows_text)
            )
            sheets_text.append(block)

    wb.close()
    return "\n\n".join(sheets_text).strip()
